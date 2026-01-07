"""
诊断图表中的文本
"""
from pptx import Presentation
import sys


def diagnose_chart(ppt_path: str, slide_num: int):
    """诊断指定幻灯片的图表"""
    prs = Presentation(ppt_path)
    
    if slide_num < 1 or slide_num > len(prs.slides):
        print(f"幻灯片 {slide_num} 不存在")
        return
    
    slide = prs.slides[slide_num - 1]
    
    print(f"幻灯片 {slide_num} 中的图表:")
    print("=" * 70)
    
    for shape_idx, shape in enumerate(slide.shapes):
        if shape.has_chart:
            print(f"\n[形状 {shape_idx}] 图表类型: {shape.chart.chart_type}")
            
            try:
                # 图表标题
                if shape.chart.has_title:
                    title = shape.chart.chart_title
                    print(f"  标题: {title.text_frame.text if title.text_frame else '无'}")
                else:
                    print(f"  标题: 无")
                
                # 分类轴
                if hasattr(shape.chart, 'category_axis'):
                    cat_axis = shape.chart.category_axis
                    print(f"  分类轴:")
                    if hasattr(cat_axis, 'axis_title') and cat_axis.axis_title:
                        print(f"    轴标题: {cat_axis.axis_title.text_frame.text if cat_axis.axis_title.text_frame else '无'}")
                    else:
                        print(f"    轴标题: 无")
                
                # 数值轴
                if hasattr(shape.chart, 'value_axis'):
                    val_axis = shape.chart.value_axis
                    print(f"  数值轴:")
                    if hasattr(val_axis, 'axis_title') and val_axis.axis_title:
                        print(f"    轴标题: {val_axis.axis_title.text_frame.text if val_axis.axis_title.text_frame else '无'}")
                    else:
                        print(f"    轴标题: 无")
                
                # 图例
                if hasattr(shape.chart, 'has_legend'):
                    print(f"  图例: {'有' if shape.chart.has_legend else '无'}")
                
                # 尝试访问图表数据
                if hasattr(shape.chart, 'plot_area'):
                    plot = shape.chart.plot_area
                    print(f"  绘图区: 存在")
                    
                    # 尝试访问系列
                    if hasattr(plot, 'series'):
                        print(f"  系列数: {len(plot.series) if plot.series else 0}")
                
            except Exception as e:
                print(f"  错误: {str(e)}")
                import traceback
                traceback.print_exc()


if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("使用方法: python3 diagnose_chart.py <ppt_file> <slide_num>")
        sys.exit(1)
    
    ppt_file = sys.argv[1]
    slide_num = int(sys.argv[2])
    
    diagnose_chart(ppt_file, slide_num)

