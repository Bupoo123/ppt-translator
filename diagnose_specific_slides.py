"""
诊断特定幻灯片的文本提取问题
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import re
import sys


def diagnose_slide(ppt_path: str, slide_numbers: list):
    """
    诊断特定幻灯片的文本提取问题
    
    Args:
        ppt_path: PPT文件路径
        slide_numbers: 要诊断的幻灯片编号列表（从1开始）
    """
    prs = Presentation(ppt_path)
    
    print("=" * 70)
    print(f"诊断幻灯片: {slide_numbers}")
    print("=" * 70)
    
    for slide_num in slide_numbers:
        if slide_num < 1 or slide_num > len(prs.slides):
            print(f"\n❌ 幻灯片 {slide_num} 不存在（总共 {len(prs.slides)} 张）")
            continue
        
        slide_idx = slide_num - 1
        slide = prs.slides[slide_idx]
        
        print(f"\n{'='*70}")
        print(f"幻灯片 {slide_num} (索引 {slide_idx}):")
        print(f"{'='*70}")
        print(f"总形状数: {len(slide.shapes)}")
        
        all_texts = []
        
        for shape_idx, shape in enumerate(slide.shapes):
            shape_type = shape.shape_type
            shape_type_name = MSO_SHAPE_TYPE(shape_type).name if hasattr(MSO_SHAPE_TYPE, '__call__') else str(shape_type)
            
            print(f"\n[形状 {shape_idx}] 类型: {shape_type_name} ({shape_type})")
            
            # 检查各种文本容器
            texts_found = []
            
            # 1. 标准文本框
            if shape.has_text_frame:
                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    text = para.text.strip()
                    if text:
                        chinese_count = len(re.findall(r'[\u4e00-\u9fff]', text))
                        chinese_ratio = chinese_count / len(text) if len(text) > 0 else 0
                        
                        texts_found.append({
                            'type': 'textbox',
                            'para_idx': para_idx,
                            'text': text,
                            'chinese_count': chinese_count,
                            'chinese_ratio': chinese_ratio,
                            'should_translate': chinese_ratio >= 0.2 and not text.isdigit()
                        })
                        print(f"  段落 {para_idx}: '{text[:60]}...'")
                        print(f"    中文: {chinese_count} 个, 比例: {chinese_ratio:.1%}, 应翻译: {texts_found[-1]['should_translate']}")
            
            # 2. 表格
            elif shape.has_table:
                print(f"  表格: {len(shape.table.rows)} 行 x {len(shape.table.columns)} 列")
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        text = cell.text.strip()
                        if text:
                            chinese_count = len(re.findall(r'[\u4e00-\u9fff]', text))
                            chinese_ratio = chinese_count / len(text) if len(text) > 0 else 0
                            
                            texts_found.append({
                                'type': 'table',
                                'row': row_idx,
                                'col': col_idx,
                                'text': text,
                                'chinese_count': chinese_count,
                                'chinese_ratio': chinese_ratio,
                                'should_translate': chinese_ratio >= 0.2 and not text.isdigit()
                            })
                            print(f"  单元格 [{row_idx},{col_idx}]: '{text[:60]}...'")
                            print(f"    中文: {chinese_count} 个, 比例: {chinese_ratio:.1%}, 应翻译: {texts_found[-1]['should_translate']}")
            
            # 3. 组合形状
            elif shape.shape_type == 6:  # GROUP
                print(f"  组合形状，包含 {len(shape.shapes)} 个子形状:")
                for sub_idx, sub_shape in enumerate(shape.shapes):
                    if sub_shape.has_text_frame:
                        for para_idx, para in enumerate(sub_shape.text_frame.paragraphs):
                            text = para.text.strip()
                            if text:
                                chinese_count = len(re.findall(r'[\u4e00-\u9fff]', text))
                                chinese_ratio = chinese_count / len(text) if len(text) > 0 else 0
                                
                                texts_found.append({
                                    'type': 'group_textbox',
                                    'sub_shape': sub_idx,
                                    'para_idx': para_idx,
                                    'text': text,
                                    'chinese_count': chinese_count,
                                    'chinese_ratio': chinese_ratio,
                                    'should_translate': chinese_ratio >= 0.2 and not text.isdigit()
                                })
                                print(f"    子形状 {sub_idx} 段落 {para_idx}: '{text[:60]}...'")
                                print(f"      中文: {chinese_count} 个, 比例: {chinese_ratio:.1%}, 应翻译: {texts_found[-1]['should_translate']}")
            
            # 4. 占位符
            elif shape.is_placeholder:
                print(f"  占位符 (类型: {shape.placeholder_format.type})")
                if shape.has_text_frame:
                    for para_idx, para in enumerate(shape.text_frame.paragraphs):
                        text = para.text.strip()
                        if text:
                            chinese_count = len(re.findall(r'[\u4e00-\u9fff]', text))
                            chinese_ratio = chinese_count / len(text) if len(text) > 0 else 0
                            
                            texts_found.append({
                                'type': 'placeholder',
                                'para_idx': para_idx,
                                'text': text,
                                'chinese_count': chinese_count,
                                'chinese_ratio': chinese_ratio,
                                'should_translate': chinese_ratio >= 0.2 and not text.isdigit()
                            })
                            print(f"  段落 {para_idx}: '{text[:60]}...'")
                            print(f"    中文: {chinese_count} 个, 比例: {chinese_ratio:.1%}, 应翻译: {texts_found[-1]['should_translate']}")
            
            # 5. 自动形状
            elif shape.shape_type == 1:  # AUTO_SHAPE
                print(f"  自动形状")
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    for para_idx, para in enumerate(shape.text_frame.paragraphs):
                        text = para.text.strip()
                        if text:
                            chinese_count = len(re.findall(r'[\u4e00-\u9fff]', text))
                            chinese_ratio = chinese_count / len(text) if len(text) > 0 else 0
                            
                            texts_found.append({
                                'type': 'autoshape',
                                'para_idx': para_idx,
                                'text': text,
                                'chinese_count': chinese_count,
                                'chinese_ratio': chinese_ratio,
                                'should_translate': chinese_ratio >= 0.2 and not text.isdigit()
                            })
                            print(f"  段落 {para_idx}: '{text[:60]}...'")
                            print(f"    中文: {chinese_count} 个, 比例: {chinese_ratio:.1%}, 应翻译: {texts_found[-1]['should_translate']}")
            
            all_texts.extend(texts_found)
        
        # 使用PPTProcessor提取，对比结果
        print(f"\n{'='*70}")
        print("使用PPTProcessor提取的结果:")
        print(f"{'='*70}")
        
        from ppt_processor import PPTProcessor
        processor = PPTProcessor(ppt_path)
        extracted = processor.extract_texts()
        
        # 找到对应幻灯片的提取结果
        slide_extracted = None
        for slide_data in extracted:
            if slide_data['slide_index'] == slide_idx:
                slide_extracted = slide_data
                break
        
        if slide_extracted:
            extracted_texts = [item['text'] for item in slide_extracted['texts']]
            print(f"提取到 {len(extracted_texts)} 个文本块:")
            for i, text in enumerate(extracted_texts, 1):
                print(f"  {i}. '{text[:60]}...'")
        else:
            print("未提取到任何文本块")
        
        # 对比分析
        print(f"\n{'='*70}")
        print("对比分析:")
        print(f"{'='*70}")
        print(f"总文本块: {len(all_texts)} 个")
        print(f"应翻译的文本块: {sum(1 for t in all_texts if t['should_translate'])} 个")
        print(f"实际提取的文本块: {len(extracted_texts) if slide_extracted else 0} 个")
        
        # 找出遗漏的文本
        if slide_extracted:
            extracted_set = set(extracted_texts)
            missed = []
            for t in all_texts:
                if t['should_translate'] and t['text'] not in extracted_set:
                    missed.append(t)
            
            if missed:
                print(f"\n⚠️  遗漏的文本块 ({len(missed)} 个):")
                for t in missed:
                    print(f"  - [{t['type']}] '{t['text']}'")
                    print(f"    中文比例: {t['chinese_ratio']:.1%}, 应翻译: {t['should_translate']}")
            else:
                print("\n✅ 所有应翻译的文本都被正确提取")


if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("使用方法: python3 diagnose_specific_slides.py <ppt_file> <slide_num1> [slide_num2] ...")
        print("示例: python3 diagnose_specific_slides.py file.pptx 9 19")
        sys.exit(1)
    
    ppt_file = sys.argv[1]
    slide_nums = [int(n) for n in sys.argv[2:]]
    
    diagnose_slide(ppt_file, slide_nums)

