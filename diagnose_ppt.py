"""
PPT诊断脚本 - 检查文本提取和更新的问题
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import re


def diagnose_ppt(ppt_path: str):
    """
    诊断PPT文件，找出所有可能的文本容器
    """
    prs = Presentation(ppt_path)
    
    print("=" * 70)
    print(f"PPT诊断报告: {ppt_path}")
    print("=" * 70)
    
    total_shapes = 0
    text_shapes = 0
    missed_shapes = []
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n幻灯片 {slide_idx + 1}:")
        print(f"  总形状数: {len(slide.shapes)}")
        
        for shape_idx, shape in enumerate(slide.shapes):
            total_shapes += 1
            shape_type = shape.shape_type
            
            # 检查各种可能的文本容器
            has_text = False
            text_content = ""
            
            # 1. 标准文本框
            if shape.has_text_frame:
                has_text = True
                text_content = shape.text_frame.text
                print(f"    [形状 {shape_idx}] 类型: {shape_type} (TEXT_FRAME)")
                text_shapes += 1
            
            # 2. 表格
            elif shape.has_table:
                has_text = True
                cell_texts = []
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            cell_texts.append(cell.text.strip())
                text_content = " | ".join(cell_texts)
                print(f"    [形状 {shape_idx}] 类型: {shape_type} (TABLE)")
                text_shapes += 1
            
            # 3. 组合形状（GroupShape）
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                print(f"    [形状 {shape_idx}] 类型: {shape_type} (GROUP - 需要递归检查)")
                # 组合形状需要递归检查
                group_texts = []
                for sub_shape in shape.shapes:
                    if sub_shape.has_text_frame:
                        group_texts.append(sub_shape.text_frame.text)
                if group_texts:
                    has_text = True
                    text_content = " | ".join(group_texts)
                    print(f"      -> 包含 {len(group_texts)} 个文本子形状")
                    text_shapes += 1
            
            # 4. 自动形状（可能包含文本）
            elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    has_text = True
                    text_content = shape.text_frame.text
                    print(f"    [形状 {shape_idx}] 类型: {shape_type} (AUTO_SHAPE with text)")
                    text_shapes += 1
                else:
                    print(f"    [形状 {shape_idx}] 类型: {shape_type} (AUTO_SHAPE without text)")
            
            # 5. 占位符（Placeholder）
            elif shape.is_placeholder:
                if shape.has_text_frame:
                    has_text = True
                    text_content = shape.text_frame.text
                    print(f"    [形状 {shape_idx}] 类型: {shape_type} (PLACEHOLDER)")
                    text_shapes += 1
            
            # 6. 其他类型
            else:
                # 尝试检查是否有文本
                try:
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        has_text = True
                        text_content = shape.text_frame.text
                        print(f"    [形状 {shape_idx}] 类型: {shape_type} (OTHER with text)")
                        text_shapes += 1
                    else:
                        print(f"    [形状 {shape_idx}] 类型: {shape_type} (OTHER - 无文本)")
                except:
                    print(f"    [形状 {shape_idx}] 类型: {shape_type} (OTHER - 无法检查)")
            
            # 如果有文本，显示内容预览
            if has_text and text_content.strip():
                preview = text_content[:50].replace('\n', ' ')
                chinese_count = len(re.findall(r'[\u4e00-\u9fff]', text_content))
                if chinese_count > 0:
                    print(f"      文本预览: {preview}... (包含 {chinese_count} 个中文字符)")
                else:
                    print(f"      文本预览: {preview}... (无中文)")
    
    print("\n" + "=" * 70)
    print("统计信息:")
    print(f"  总形状数: {total_shapes}")
    print(f"  包含文本的形状数: {text_shapes}")
    print(f"  可能遗漏的形状数: {total_shapes - text_shapes}")
    print("=" * 70)


def check_text_extraction_issues(ppt_path: str):
    """
    检查文本提取的具体问题
    """
    prs = Presentation(ppt_path)
    
    print("\n" + "=" * 70)
    print("文本提取问题检查")
    print("=" * 70)
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n幻灯片 {slide_idx + 1}:")
        
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                paragraphs = shape.text_frame.paragraphs
                print(f"  形状 {shape_idx}: {len(paragraphs)} 个段落")
                
                for para_idx, para in enumerate(paragraphs):
                    text = para.text.strip()
                    if text:
                        runs_count = len(para.runs)
                        print(f"    段落 {para_idx}: '{text[:30]}...' ({runs_count} 个runs)")
                        
                        # 检查是否有多个runs（可能导致更新问题）
                        if runs_count > 1:
                            print(f"      ⚠️  警告: 此段落包含多个runs，更新时可能需要特殊处理")


if __name__ == '__main__':
    import sys
    
    if len(sys.argv) < 2:
        print("使用方法: python3 diagnose_ppt.py <ppt_file>")
        sys.exit(1)
    
    ppt_file = sys.argv[1]
    
    diagnose_ppt(ppt_file)
    check_text_extraction_issues(ppt_file)

