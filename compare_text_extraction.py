"""
对比文本提取 - 检查哪些文本被识别，哪些被遗漏
"""
from pptx import Presentation
import re
import sys
import os


def extract_all_texts(ppt_path: str):
    """
    提取PPT中所有可能的文本（不判断是否需要翻译）
    """
    prs = Presentation(ppt_path)
    all_texts = []
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_texts = []
        
        for shape_idx, shape in enumerate(slide.shapes):
            # 标准文本框
            if shape.has_text_frame:
                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    text = para.text.strip()
                    if text:
                        chinese_count = len(re.findall(r'[\u4e00-\u9fff]', text))
                        slide_texts.append({
                            'shape_index': shape_idx,
                            'shape_type': 'textbox',
                            'paragraph_index': para_idx,
                            'text': text,
                            'chinese_count': chinese_count,
                            'length': len(text)
                        })
            
            # 表格
            elif shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        text = cell.text.strip()
                        if text:
                            chinese_count = len(re.findall(r'[\u4e00-\u9fff]', text))
                            slide_texts.append({
                                'shape_index': shape_idx,
                                'shape_type': 'table',
                                'row_index': row_idx,
                                'col_index': col_idx,
                                'text': text,
                                'chinese_count': chinese_count,
                                'length': len(text)
                            })
            
            # 组合形状
            elif shape.shape_type == 6:  # GROUP
                for sub_idx, sub_shape in enumerate(shape.shapes):
                    if sub_shape.has_text_frame:
                        for para_idx, para in enumerate(sub_shape.text_frame.paragraphs):
                            text = para.text.strip()
                            if text:
                                chinese_count = len(re.findall(r'[\u4e00-\u9fff]', text))
                                slide_texts.append({
                                    'shape_index': shape_idx,
                                    'shape_type': 'group_textbox',
                                    'sub_shape_index': sub_idx,
                                    'paragraph_index': para_idx,
                                    'text': text,
                                    'chinese_count': chinese_count,
                                    'length': len(text)
                                })
        
        if slide_texts:
            all_texts.append({
                'slide_index': slide_idx,
                'texts': slide_texts
            })
    
    return all_texts


def compare_extraction(ppt_path: str):
    """
    对比文本提取结果
    """
    print("=" * 70)
    print(f"文本提取对比分析: {ppt_path}")
    print("=" * 70)
    
    # 提取所有文本
    all_texts = extract_all_texts(ppt_path)
    
    # 使用PPTProcessor提取（会过滤不需要翻译的）
    from ppt_processor import PPTProcessor
    processor = PPTProcessor(ppt_path)
    filtered_texts = processor.extract_texts()
    
    print(f"\n统计信息:")
    print(f"  总幻灯片数: {len(all_texts)} 张")
    print(f"  总文本块数: {sum(len(s['texts']) for s in all_texts)} 个")
    print(f"  识别为可翻译的幻灯片: {len(filtered_texts)} 张")
    print(f"  识别为可翻译的文本块: {sum(len(s['texts']) for s in filtered_texts)} 个")
    
    # 创建过滤后的文本集合（用于对比）
    filtered_text_set = set()
    for slide_data in filtered_texts:
        for item in slide_data['texts']:
            filtered_text_set.add(item['text'].strip())
    
    # 找出被过滤掉的文本
    missed_texts = []
    for slide_data in all_texts:
        for item in slide_data['texts']:
            text = item['text'].strip()
            if text not in filtered_text_set:
                # 检查为什么被过滤
                chinese_ratio = item['chinese_count'] / item['length'] if item['length'] > 0 else 0
                missed_texts.append({
                    'slide_index': slide_data['slide_index'],
                    'text': text,
                    'chinese_count': item['chinese_count'],
                    'chinese_ratio': chinese_ratio,
                    'length': item['length'],
                    'shape_type': item['shape_type']
                })
    
    if missed_texts:
        print(f"\n⚠️  发现 {len(missed_texts)} 个文本块被过滤（未识别为可翻译）:")
        print("\n可能的原因:")
        
        # 按原因分类
        no_chinese = [t for t in missed_texts if t['chinese_count'] == 0]
        low_ratio = [t for t in missed_texts if t['chinese_count'] > 0 and t['chinese_ratio'] < 0.2]
        digits_only = [t for t in missed_texts if t['text'].isdigit()]
        
        if digits_only:
            print(f"  1. 纯数字 ({len(digits_only)} 个):")
            for t in digits_only[:5]:
                print(f"     幻灯片{t['slide_index']+1}: '{t['text']}'")
        
        if no_chinese:
            print(f"  2. 无中文 ({len(no_chinese)} 个):")
            for t in no_chinese[:10]:
                print(f"     幻灯片{t['slide_index']+1}: '{t['text'][:50]}...'")
        
        if low_ratio:
            print(f"  3. 中文比例<20% ({len(low_ratio)} 个):")
            for t in low_ratio[:10]:
                print(f"     幻灯片{t['slide_index']+1}: '{t['text'][:50]}...' (中文比例: {t['chinese_ratio']:.1%})")
        
        # 显示前20个被过滤的文本
        print(f"\n前20个被过滤的文本示例:")
        for i, t in enumerate(missed_texts[:20], 1):
            print(f"  {i}. 幻灯片{t['slide_index']+1} [{t['shape_type']}]: '{t['text'][:60]}...'")
    else:
        print("\n✅ 所有文本都被正确识别！")
    
    # 显示每张幻灯片的详细信息
    print(f"\n详细分析（前10张幻灯片）:")
    for slide_data in all_texts[:10]:
        slide_idx = slide_data['slide_index']
        total = len(slide_data['texts'])
        
        # 统计被过滤的数量
        filtered_count = 0
        for item in slide_data['texts']:
            if item['text'].strip() in filtered_text_set:
                filtered_count += 1
        
        print(f"\n幻灯片 {slide_idx + 1}:")
        print(f"  总文本块: {total} 个")
        print(f"  识别为可翻译: {filtered_count} 个")
        print(f"  被过滤: {total - filtered_count} 个")


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("使用方法: python3 compare_text_extraction.py <ppt_file>")
        sys.exit(1)
    
    ppt_file = sys.argv[1]
    compare_extraction(ppt_file)

