"""
PPT处理器 - 核心功能模块
负责PPT的解析、文本提取和回填
"""
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import List, Dict, Tuple, Optional
import re


class PPTProcessor:
    """PPT处理器类"""
    
    def __init__(self, ppt_path: str):
        """
        初始化PPT处理器
        
        Args:
            ppt_path: PPT文件路径
        """
        self.ppt_path = ppt_path
        self.prs = Presentation(ppt_path)
        self.slides_data = []
    
    def extract_texts(self) -> List[Dict]:
        """
        提取所有可翻译的文本
        
        Returns:
            包含文本信息的列表，每个元素包含：
            - slide_index: 幻灯片索引
            - shape_index: 形状索引
            - text: 原始文本
            - text_type: 文本类型（textbox, table, chart等）
        """
        texts = []
        
        for slide_idx, slide in enumerate(self.prs.slides):
            slide_texts = []
            
            for shape_idx, shape in enumerate(slide.shapes):
                # 处理文本框（包括占位符）
                if shape.has_text_frame:
                    for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                        text = paragraph.text.strip()
                        if text and self._should_translate(text):
                            slide_texts.append({
                                'slide_index': slide_idx,
                                'shape_index': shape_idx,
                                'paragraph_index': para_idx,  # 修复：使用实际的段落索引
                                'text': text,
                                'text_type': 'textbox',
                                'shape': shape,
                                'paragraph': paragraph
                            })
                
                # 处理组合形状（GroupShape）中的文本
                elif shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                    for sub_shape_idx, sub_shape in enumerate(shape.shapes):
                        if sub_shape.has_text_frame:
                            for para_idx, paragraph in enumerate(sub_shape.text_frame.paragraphs):
                                text = paragraph.text.strip()
                                if text and self._should_translate(text):
                                    slide_texts.append({
                                        'slide_index': slide_idx,
                                        'shape_index': shape_idx,
                                        'sub_shape_index': sub_shape_idx,
                                        'paragraph_index': para_idx,
                                        'text': text,
                                        'text_type': 'group_textbox',
                                        'shape': shape,
                                        'sub_shape': sub_shape,
                                        'paragraph': paragraph
                                    })
                
                # 处理表格
                if shape.has_table:
                    for row_idx, row in enumerate(shape.table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            text = cell.text.strip()
                            if text and self._should_translate(text):
                                slide_texts.append({
                                    'slide_index': slide_idx,
                                    'shape_index': shape_idx,
                                    'row_index': row_idx,
                                    'col_index': col_idx,
                                    'text': text,
                                    'text_type': 'table',
                                    'cell': cell
                                })
                
                # 处理图表（Chart）中的文本
                if shape.has_chart:
                    try:
                        chart = shape.chart
                        
                        # 图表标题
                        if chart.has_title and chart.chart_title:
                            title_text = chart.chart_title.text_frame.text.strip()
                            if title_text and self._should_translate(title_text):
                                slide_texts.append({
                                    'slide_index': slide_idx,
                                    'shape_index': shape_idx,
                                    'text': title_text,
                                    'text_type': 'chart_title',
                                    'chart': chart
                                })
                        
                        # 坐标轴标签
                        if hasattr(chart, 'category_axis') and chart.category_axis:
                            if hasattr(chart.category_axis, 'axis_title') and chart.category_axis.axis_title:
                                axis_text = chart.category_axis.axis_title.text_frame.text.strip()
                                if axis_text and self._should_translate(axis_text):
                                    slide_texts.append({
                                        'slide_index': slide_idx,
                                        'shape_index': shape_idx,
                                        'text': axis_text,
                                        'text_type': 'chart_axis',
                                        'axis_type': 'category',
                                        'chart': chart
                                    })
                        
                        if hasattr(chart, 'value_axis') and chart.value_axis:
                            if hasattr(chart.value_axis, 'axis_title') and chart.value_axis.axis_title:
                                axis_text = chart.value_axis.axis_title.text_frame.text.strip()
                                if axis_text and self._should_translate(axis_text):
                                    slide_texts.append({
                                        'slide_index': slide_idx,
                                        'shape_index': shape_idx,
                                        'text': axis_text,
                                        'text_type': 'chart_axis',
                                        'axis_type': 'value',
                                        'chart': chart
                                    })
                        
                        # 图例（如果可访问）
                        if hasattr(chart, 'legend') and chart.legend:
                            if hasattr(chart.legend, 'text_frame') and chart.legend.text_frame:
                                legend_text = chart.legend.text_frame.text.strip()
                                if legend_text and self._should_translate(legend_text):
                                    slide_texts.append({
                                        'slide_index': slide_idx,
                                        'shape_index': shape_idx,
                                        'text': legend_text,
                                        'text_type': 'chart_legend',
                                        'chart': chart
                                    })
                    except Exception as e:
                        # 图表处理可能失败，忽略错误继续处理其他形状
                        pass
            
            if slide_texts:
                texts.append({
                    'slide_index': slide_idx,
                    'texts': slide_texts
                })
        
        self.slides_data = texts
        return texts
    
    def _should_translate(self, text: str) -> bool:
        """
        判断文本是否需要翻译
        
        Args:
            text: 待判断的文本
            
        Returns:
            True表示需要翻译，False表示跳过
        """
        # 跳过纯数字
        if text.isdigit():
            return False
        
        # 检查是否包含中文
        chinese_chars = len(re.findall(r'[\u4e00-\u9fff]', text))
        if chinese_chars == 0:
            return False
        
        total_chars = len(text)
        if total_chars > 0:
            chinese_ratio = chinese_chars / total_chars
            
            # 如果包含中文，但比例较低，需要更宽松的判断
            # 对于包含冒号、时间单位等的情况，降低阈值
            if chinese_ratio < 0.2:
                # 特殊处理：如果文本包含常见的中文标点或单位，即使比例低也翻译
                has_chinese_punctuation = bool(re.search(r'[：，。、；]', text))
                has_time_unit = bool(re.search(r'[小时|分钟|秒|天|年|月]', text))
                has_chinese_keywords = bool(re.search(r'[轮次|金额|融资|投资|成本|价格|数量]', text))
                
                # 如果包含这些特征，即使比例低于20%也翻译
                if has_chinese_punctuation or has_time_unit or has_chinese_keywords:
                    return True
                
                return False
        
        # 包含中文的文本需要翻译
        return True
    
    def _preserve_format_and_set_font(self, paragraph, translated_text: str):
        """
        保留原有格式并设置字体为Arial
        
        Args:
            paragraph: 段落对象
            translated_text: 翻译后的文本
        """
        # 保存原有格式（从第一个run获取，如果没有run则使用默认值）
        font_size = None
        font_color = None
        font_bold = False
        font_italic = False
        east_asian_font = None  # 中文字体
        
        if paragraph.runs:
            first_run = paragraph.runs[0]
            if first_run.font.size:
                font_size = first_run.font.size
            if first_run.font.color and first_run.font.color.type:
                if first_run.font.color.type == 1:  # RGB color
                    font_color = first_run.font.color.rgb
            font_bold = first_run.font.bold if first_run.font.bold is not None else False
            font_italic = first_run.font.italic if first_run.font.italic is not None else False
            # 尝试获取中文字体（如果有）
            if hasattr(first_run.font, 'name') and first_run.font.name:
                # 保存原字体名称，可能包含中文字体信息
                pass
        
        # 清除原有内容
        paragraph.clear()
        
        # 检查是否包含中文字符
        has_chinese = bool(re.search(r'[\u4e00-\u9fff]', translated_text))
        
        # 添加新文本
        run = paragraph.add_run()
        run.text = translated_text
        
        # 应用保存的格式
        if font_size:
            run.font.size = font_size
        if font_color:
            run.font.color.rgb = font_color
        run.font.bold = font_bold
        run.font.italic = font_italic
        
        # 设置字体：英文字体用Arial
        # 如果包含中文，中文字体保持原样（不设置，使用系统默认或原字体）
        run.font.name = 'Arial'
        
        # 注意：python-pptx的font.name主要影响拉丁字符
        # 中文字符通常使用系统默认字体或通过其他方式设置
        # 这里我们只设置英文字体为Arial，中文字体保持原样
    
    def update_text(self, slide_index: int, shape_index: int, 
                   original_text: str, translated_text: str,
                   paragraph_index: int = None, 
                   row_index: int = None, col_index: int = None,
                   sub_shape_index: int = None, **kwargs):
        """
        更新PPT中的文本
        
        Args:
            slide_index: 幻灯片索引
            shape_index: 形状索引
            original_text: 原始文本
            translated_text: 翻译后的文本
            paragraph_index: 段落索引（用于textbox）
            row_index: 行索引（用于table）
            col_index: 列索引（用于table）
            sub_shape_index: 子形状索引（用于组合形状）
        """
        slide = self.prs.slides[slide_index]
        shape = slide.shapes[shape_index]
        
        # 处理组合形状中的文本
        if sub_shape_index is not None and shape.shape_type == 6:  # GROUP
            sub_shape = shape.shapes[sub_shape_index]
            if sub_shape.has_text_frame:
                if paragraph_index is not None and paragraph_index < len(sub_shape.text_frame.paragraphs):
                    paragraph = sub_shape.text_frame.paragraphs[paragraph_index]
                    # 保留格式并设置字体
                    self._preserve_format_and_set_font(paragraph, translated_text)
                else:
                    # 尝试找到匹配的段落
                    for para in sub_shape.text_frame.paragraphs:
                        if original_text.strip() in para.text or para.text.strip() == original_text.strip():
                            self._preserve_format_and_set_font(para, translated_text)
                            break
        
        elif shape.has_text_frame:
            # 更新文本框
            if paragraph_index is not None and paragraph_index < len(shape.text_frame.paragraphs):
                # 如果指定了段落索引且有效，更新该段落
                paragraph = shape.text_frame.paragraphs[paragraph_index]
                # 保留格式并设置字体
                self._preserve_format_and_set_font(paragraph, translated_text)
            else:
                # 如果没有指定段落或索引无效，尝试找到包含原始文本的段落
                found = False
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    # 精确匹配或包含匹配
                    if para_text == original_text.strip() or original_text.strip() in para_text:
                        self._preserve_format_and_set_font(paragraph, translated_text)
                        found = True
                        break
                
                # 如果找不到匹配的段落，尝试更新整个文本框（但保留其他段落）
                if not found:
                    # 只更新包含原始文本的段落，而不是整个文本框
                    for paragraph in shape.text_frame.paragraphs:
                        if original_text.strip() in paragraph.text:
                            self._preserve_format_and_set_font(paragraph, translated_text)
                            break
        
        elif shape.has_table and row_index is not None and col_index is not None:
            # 更新表格单元格
            cell = shape.table.rows[row_index].cells[col_index]
            # 保存单元格格式
            if cell.text_frame.paragraphs:
                paragraph = cell.text_frame.paragraphs[0]
                self._preserve_format_and_set_font(paragraph, translated_text)
            else:
                cell.text = translated_text
                # 设置字体
                if cell.text_frame.paragraphs:
                    for run in cell.text_frame.paragraphs[0].runs:
                        run.font.name = 'Arial'
        
        elif shape.has_chart:
            # 更新图表文本
            try:
                chart = shape.chart
                text_type = kwargs.get('text_type', '')
                
                if text_type == 'chart_title' and chart.has_title:
                    chart.chart_title.text_frame.text = translated_text
                elif text_type == 'chart_axis':
                    axis_type = kwargs.get('axis_type', 'category')
                    if axis_type == 'category' and hasattr(chart, 'category_axis'):
                        if hasattr(chart.category_axis, 'axis_title'):
                            chart.category_axis.axis_title.text_frame.text = translated_text
                    elif axis_type == 'value' and hasattr(chart, 'value_axis'):
                        if hasattr(chart.value_axis, 'axis_title'):
                            chart.value_axis.axis_title.text_frame.text = translated_text
                elif text_type == 'chart_legend' and hasattr(chart, 'legend'):
                    if hasattr(chart.legend, 'text_frame'):
                        chart.legend.text_frame.text = translated_text
            except Exception as e:
                # 图表更新可能失败，记录错误但继续
                pass
    
    def save(self, output_path: str):
        """
        保存PPT文件
        
        Args:
            output_path: 输出文件路径
        """
        self.prs.save(output_path)
    
    def get_slide_texts(self, slide_index: int) -> List[str]:
        """
        获取指定幻灯片的所有文本（用于上下文翻译）
        
        Args:
            slide_index: 幻灯片索引
            
        Returns:
            文本列表
        """
        texts = []
        slide = self.prs.slides[slide_index]
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text and self._should_translate(text):
                        texts.append(text)
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text and self._should_translate(text):
                            texts.append(text)
            elif shape.shape_type == 6:  # GROUP - 处理组合形状
                for sub_shape in shape.shapes:
                    if sub_shape.has_text_frame:
                        for paragraph in sub_shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text and self._should_translate(text):
                                texts.append(text)
        
        return texts

